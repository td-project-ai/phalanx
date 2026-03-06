"""
PPTX Renderer — converts a parsed outline + theme into a PowerPoint file.

Usage:
    python tools/presentations/render_pptx.py outline.yaml [-o output.pptx] [-t credera]

    from tools.presentations.render_pptx import render
    render("outline.yaml", output="deck.pptx", theme="credera")

Architecture:
    1. Parse outline YAML → normalised slide dicts (via parse_outline.py)
    2. Load theme config.yaml → colour tokens, font config, layout mapping
    3. Open template.pptx → remove sample slides, keep masters/layouts
    4. For each slide in outline:
       a. Pick layout from theme mapping
       b. Add slide with that layout
       c. Dispatch to component renderer
    5. Save .pptx
"""

from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt

# Local imports
from .parse_outline import parse as parse_outline
from .components.typography import (
    render_slide_chrome,
    render_title_slide,
    render_section_slide,
    render_stat_slide,
    render_quote_slide,
)
from .components.layouts import render_content, render_two_col
from .components.data import (
    render_stat_grid,
    render_card_grid,
    render_data_table,
    render_highlight_grid,
)
from .components.flow import (
    render_step_flow,
    render_funnel,
    render_timeline,
    render_ascend,
)
from .components.structure import (
    render_layer_stack,
    render_hub_spoke,
    render_process_loop,
    render_comparison,
)

# ---------------------------------------------------------------------------
# Theme loader
# ---------------------------------------------------------------------------

BRANDS_DIR = Path(__file__).resolve().parent.parent.parent / "context" / "templates" / "presentations" / "themes" / "brands"


def _extract_pptx_config(brand_file: Path) -> str | None:
    """Extract PPTX config YAML from a brand .md file.

    Looks for a ``## PPTX Config`` heading followed by a fenced yaml block.
    Returns the YAML string or None if not found.
    """
    text = brand_file.read_text(encoding="utf-8")
    in_section = False
    in_code = False
    yaml_lines: list[str] = []

    for line in text.splitlines():
        if line.startswith("## PPTX Config"):
            in_section = True
            continue
        if in_section and line.startswith("## "):
            break  # hit next section
        if in_section and line.strip().startswith("```yaml"):
            in_code = True
            continue
        if in_code and line.strip() == "```":
            break
        if in_code:
            yaml_lines.append(line)

    return "\n".join(yaml_lines) if yaml_lines else None


def load_theme(theme_name: str, variant: str = "default") -> dict:
    """Load a PPTX theme config and apply variant overrides.

    Resolution order:
      1. Brand .md file with embedded ``## PPTX Config`` YAML block
      2. Standalone ``config.yaml`` in brands/<name>/ (legacy fallback)

    The binary ``template.pptx`` lives at brands/<name>/template.pptx.
    """
    theme: dict | None = None

    # 1. Try brand .md file first
    brand_file = BRANDS_DIR / f"{theme_name}.md"
    if brand_file.exists():
        config_yaml = _extract_pptx_config(brand_file)
        if config_yaml:
            theme = yaml.safe_load(config_yaml)

    # 2. Fallback to standalone config.yaml in brand subdirectory
    if theme is None:
        brand_dir = BRANDS_DIR / theme_name
        config_path = brand_dir / "config.yaml"
        if not config_path.exists():
            raise FileNotFoundError(
                f"No PPTX config found for '{theme_name}'. "
                f"Checked: {brand_file}, {config_path}"
            )
        with open(config_path, "r", encoding="utf-8") as f:
            theme = yaml.safe_load(f)

    # Apply variant overrides
    if variant and variant != "default":
        variants = theme.get("variants", {})
        if variant in variants:
            overrides = variants[variant]
            theme["colors"] = {**theme.get("colors", {}), **overrides}

    # Template.pptx lives in the brand subdirectory
    brand_dir = BRANDS_DIR / theme_name
    theme["_template_path"] = str(brand_dir / "template.pptx")
    theme["_theme_dir"] = str(brand_dir)

    return theme


# ---------------------------------------------------------------------------
# Layout resolution
# ---------------------------------------------------------------------------

def get_layout(prs: Presentation, theme: dict, layout_key: str,
               slide_data: dict | None = None):
    """Get a slide layout from the presentation by theme mapping.

    If the theme maps the key to "draw", returns the blank layout.
    Otherwise returns the layout at the specified index.
    Handles compound keys like section → section_gradient / section_photo.
    """
    layouts_map = theme.get("layouts", {})

    # Handle section style variants
    if layout_key == "section" and slide_data:
        style = slide_data.get("style", "gradient")
        compound_key = f"section_{style}"
        if compound_key in layouts_map:
            layout_key = compound_key

    layout_val = layouts_map.get(layout_key, "draw")

    if layout_val == "draw":
        # Use blank layout (index 0)
        blank_idx = layouts_map.get("blank", 0)
        return prs.slide_layouts[blank_idx]
    elif isinstance(layout_val, int):
        return prs.slide_layouts[layout_val]
    else:
        # Fallback to blank
        return prs.slide_layouts[0]


# ---------------------------------------------------------------------------
# Slide dispatch
# ---------------------------------------------------------------------------

RENDERERS = {
    # Structural (use native template layouts)
    "title": render_title_slide,
    "section": render_section_slide,
    "stat": render_stat_slide,
    "quote": render_quote_slide,
    # end: no renderer, just the layout

    # Content
    "content": render_content,
    "two-col": render_two_col,

    # Data
    "stat-grid": render_stat_grid,
    "card-grid": render_card_grid,
    "data-table": render_data_table,
    "highlight-grid": render_highlight_grid,

    # Flow
    "step-flow": render_step_flow,
    "funnel": render_funnel,
    "timeline": render_timeline,
    "ascend": render_ascend,

    # Structure
    "layer-stack": render_layer_stack,
    "hub-spoke": render_hub_spoke,
    "process-loop": render_process_loop,
    "comparison": render_comparison,
}


def render_slide(prs: Presentation, slide_data: dict, theme: dict):
    """Add a slide and dispatch to the appropriate renderer."""
    layout_key = slide_data.get("layout", "content")
    layout = get_layout(prs, theme, layout_key, slide_data)
    slide = prs.slides.add_slide(layout)

    renderer = RENDERERS.get(layout_key)
    if renderer:
        renderer(slide, slide_data, theme)
    elif layout_key == "end":
        # End splash uses template layout, add speaker notes if provided
        if "notes" in slide_data:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data["notes"]
    elif layout_key == "blank":
        # Blank — render chrome only if there's a title
        if "title" in slide_data:
            render_slide_chrome(slide, slide_data, theme)

    return slide


# ---------------------------------------------------------------------------
# Main render pipeline
# ---------------------------------------------------------------------------

def _remove_existing_slides(prs: Presentation):
    """Remove all existing slides from the presentation (keep layouts/masters)."""
    # python-pptx doesn't have a direct delete_slide method
    # We need to manipulate the XML
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId is None:
            # Try the 'id' attribute directly
            slide_id_el = prs.slides._sldIdLst[0]
            # Remove from relationship
            prs.part.drop_rel(prs.slides._sldIdLst[0].get("r:id") or "")
            prs.slides._sldIdLst.remove(slide_id_el)
        else:
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def render(outline_path: str | Path, output: str | Path | None = None,
           theme: str | None = None, variant: str | None = None) -> Path:
    """Render a PPTX from an outline YAML file.

    Args:
        outline_path: Path to outline YAML
        output: Output .pptx path (default: same name as outline with .pptx)
        theme: Theme name override (default: from outline meta)
        variant: Variant override (default: from outline meta)

    Returns:
        Path to the generated .pptx file
    """
    # 1. Parse outline
    deck = parse_outline(outline_path)
    meta = deck["meta"]
    slides_data = deck["slides"]

    # 2. Load theme
    theme_name = theme or meta.get("theme", "credera")
    variant_name = variant or meta.get("variant", "default")
    theme_cfg = load_theme(theme_name, variant_name)

    # 3. Open template
    template_path = theme_cfg["_template_path"]
    prs = Presentation(template_path)

    # 4. Remove existing slides
    _remove_existing_slides(prs)

    # 5. Set metadata
    if meta.get("title"):
        prs.core_properties.title = meta["title"]
    if meta.get("author"):
        prs.core_properties.author = meta["author"]
    date_str = meta.get("date", "")
    if not date_str:
        prs.core_properties.last_modified_by = f"PPTX Engine ({date.today().isoformat()})"

    # 6. Render each slide
    for slide_data in slides_data:
        render_slide(prs, slide_data, theme_cfg)

    # 7. Save output
    if output is None:
        outline_p = Path(outline_path)
        output = outline_p.with_suffix(".pptx")
    else:
        output = Path(output)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    return output


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="Render PPTX from outline YAML")
    ap.add_argument("outline", help="Path to outline YAML file")
    ap.add_argument("-o", "--output", help="Output .pptx path")
    ap.add_argument("-t", "--theme", help="Theme name (overrides outline meta)")
    ap.add_argument("-v", "--variant", help="Colour variant (overrides outline meta)")
    args = ap.parse_args()

    out = render(args.outline, args.output, args.theme, args.variant)
    print(f"Generated: {out} ({len(parse_outline(args.outline)['slides'])} slides)", file=sys.stderr)


if __name__ == "__main__":
    main()
